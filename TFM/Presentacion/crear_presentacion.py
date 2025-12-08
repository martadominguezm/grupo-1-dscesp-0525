#!/usr/bin/env python3
"""
Script para crear presentación PowerPoint del TFM easyMoney
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_items):
    """Slide con título y bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

    return slide

def add_section_title(prs, title, subtitle=""):
    """Slide separador de sección"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

# ==========================================
# TAREA 1: CONTEXTO + EDA + POWERBI
# ==========================================

# Slide 1: PORTADA
add_title_slide(
    prs,
    "easyMoney: Maximización de Rentabilidad mediante Analítica Predictiva",
    "Estrategia de Marketing Data-Driven para Optimizar Captación y Retención\nGrupo 1 - Nuclio Digital School 2025"
)

# Slide 2: SITUACIÓN ACTUAL
slide = add_content_slide(
    prs,
    "easyMoney: Crisis de Rentabilidad",
    [
        "Empresa: Plataforma multi-canal de productos financieros (4 años)",
        "Problema crítico: Fondos agotados sin EBITDA positivo",
        "Presión inversores: Lion Global exige rentabilidad AHORA",
        "Cambio estratégico: De captación agresiva → Rentabilización cartera actual",
        "",
        "Cifras clave:",
        "• 442,000 clientes activos (Mayo 2019)",
        "• 5 familias de productos",
        "• Desafío: ¿Cómo maximizar ingresos sin invertir en captación masiva?"
    ]
)

# Slide 3: NUESTRA SOLUCIÓN
slide = add_content_slide(
    prs,
    "Metodología: De Datos a Decisiones Comerciales",
    [
        "Los 4 pilares:",
        "",
        "1. 📊 EDA + CALIDAD DE DATOS (PowerBI)",
        "   Exploración 442K clientes, 17 meses históricos, Dashboard ejecutivo",
        "",
        "2. 🎯 PROPENSIÓN DE COMPRA",
        "   Modelos predictivos por producto, Priorización de campañas",
        "",
        "3. 👥 SEGMENTACIÓN (CLUSTERING)",
        "   6 perfiles de cliente, Personalización de mensajes",
        "",
        "4. 💰 VALOR DEL CLIENTE",
        "   Priorización por margen, ROI de campañas"
    ]
)

# Slide 4: POWERBI - VISIÓN 360º
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
slide.shapes.title.text = "Dashboard Ejecutivo: Monitorización en Tiempo Real"

# Agregar textbox con métricas
left = Inches(0.5)
top = Inches(1.5)
width = Inches(3)
height = Inches(4)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "Métricas Mayo 2019:\n\n• 10,211 ventas\n• €8.07M beneficio\n• €606.88 margen medio\n\nInsights:\n• Pico Octubre: 28K ventas\n• Caída vs 2018: -3K/mes\n• Madrid/Málaga lideran"
tf.paragraphs[0].font.size = Pt(16)

# Placeholder para imagen
left = Inches(4)
top = Inches(1.5)
width = Inches(5.5)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "[INSERTAR IMAGEN 3:\nDashboard temporal y geográfico]"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 5: POWERBI - PERFIL DEMOGRÁFICO
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "¿Dónde Están Nuestros Mejores Clientes?"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(3)
height = Inches(4)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "Perfil dominante:\n\n• 67% tienen 20-30 años\n• 296K clientes millennials\n\nGeografía:\n• País Vasco: 0.76 ratio\n• Cantabria: 0.70\n• Madrid: 0.63\n\nInsight:\nRegiones pequeñas convierten mejor"
tf.paragraphs[0].font.size = Pt(16)

left = Inches(4)
top = Inches(1.5)
width = Inches(5.5)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "[INSERTAR IMAGEN 4:\nAnálisis demográfico y territorial]"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 6: DISTRIBUCIÓN MÁRGENES
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Anatomía de la Rentabilidad por Producto"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(3.5)
height = Inches(4.5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "Dos mundos diferentes:\n\n🔵 BAJO MARGEN (€20-100):\n• Cuentas: €40-80\n• Tarjetas: €40-80\n• Lead generation\n\n🔵 ALTO MARGEN (€1K-15K):\n• Inversión: €1K-10K\n• Pensión: hasta €15K\n• Préstamo: €1K-10K\n• Upselling prioritario\n\nConclusión:\n1 venta Pensión = 100-250 ventas Tarjeta"
tf.paragraphs[0].font.size = Pt(14)

left = Inches(4.5)
top = Inches(1.5)
width = Inches(5)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "[INSERTAR IMAGEN 1:\nMárgenes por producto]"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Slide 7: LA GALLINA DE LOS HUEVOS DE ORO
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "¿Qué Productos Nos Dan De Comer?"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(3.5)
height = Inches(4.5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "🏆 PENSIÓN DOMINA:\n\n• 79% del margen total\n• €5,976 margen promedio\n• Solo 19,369 clientes (4.4%)\n\n💰 OPORTUNIDAD:\n\nSi subimos penetración a 10%:\n• +22,000 conversiones\n• +€132M revenue potencial\n\nDecisión estratégica:\nPrioridad #1 = Planes de Pensiones"
tf.paragraphs[0].font.size = Pt(14)

left = Inches(4.5)
top = Inches(1.5)
width = Inches(5)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "[INSERTAR IMAGEN 5:\nMargen por familia + Pie chart]"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ==========================================
# TAREA 2: MODELOS DE PROPENSIÓN
# ==========================================

add_section_title(prs, "TAREA 2: MODELOS DE PROPENSIÓN", "Del 'Spray & Pray' al Marketing de Precisión")

# Slide 8: ¿QUÉ ES PROPENSIÓN?
slide = add_content_slide(
    prs,
    "Del 'Spray & Pray' al Marketing de Precisión",
    [
        "Problema tradicional:",
        "❌ Campañas masivas a toda la base → Conversión <1%",
        "❌ ROI negativo, saturación, unsubscribes",
        "",
        "Solución: Propensión de Compra",
        "✅ Predecir quién comprará ANTES de lanzar campaña",
        "✅ Priorizar clientes con mayor probabilidad",
        "✅ Personalizar mensajes por perfil",
        "✅ Maximizar conversión y ROI",
        "",
        "Variables clave: Recencia, Frecuencia, Comportamiento,",
        "Sociodemográficos, Canal de entrada"
    ]
)

# Slide 9: CONSTRUCCIÓN MODELOS
slide = add_content_slide(
    prs,
    "Tres Modelos, Tres Oportunidades",
    [
        "1. PENSION PLAN (Margen: €6,071/cliente)",
        "   Target: 425K clientes elegibles",
        "",
        "2. PAYMENT CARD (Margen: €60/cliente)",
        "   Target: 440K clientes elegibles",
        "",
        "3. INVESTMENT FAMILY (Margen: €1,699/cliente)",
        "   Target: 431K clientes elegibles",
        "",
        "Técnica: XGBoost (Machine Learning)",
        "Datos: 17 meses históricos (2018-2019)",
        "Métricas: Recall 90-97%, ROC AUC 0.98-0.99"
    ]
)

# Slide 10: SEGMENTOS PROPENSIÓN
slide = add_content_slide(
    prs,
    "De 442K Clientes a Listas Accionables",
    [
        "Segmentación por propensión (score 0-1):",
        "",
        "🔥 TOP (≥0.80): Campaña premium, incentivo alto, contacto directo",
        "🔥 HIGH (0.60-0.80): Campaña estándar, incentivo medio",
        "⚠️ MEDIUM (0.50-0.60): Campaña soft, remarketing",
        "❌ MID-LOW (0.40-0.50): Nurturing, sin inversión",
        "❌ LOW (<0.40): No contactar (coste > beneficio)",
        "",
        "PENSION PLAN: 14,403 clientes Top+High (3.4% base)",
        "PAYMENT CARD: 13,891 clientes Top+High",
        "INVESTMENT: 12,184 clientes Top+High"
    ]
)

# Slide 11: VALIDACIÓN
slide = add_content_slide(
    prs,
    "¿Funciona en la Vida Real? SÍ",
    [
        "Validación en datos reales Abril 2019:",
        "",
        "Sin modelo (orgánico):",
        "• Base: 85,232 clientes → 146 conversiones",
        "• Tasa conversión: 0.17%",
        "",
        "Con modelo (Top+High):",
        "• Base: 13,000 clientes → 502 conversiones",
        "• Tasa conversión: 3.8%",
        "",
        "🎯 MEJORA: 22x más conversión (3.8% vs 0.17%)",
        "🎯 Triplica captación mensual con solo 15% contactos",
        "",
        "Conversión por producto:",
        "• Pension Plan: 3.8% | Payment Card: 8.9% | Investment: 0.19%"
    ]
)

# Slide 12: PROPENSIÓN × CLUSTERING
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Propensión × Clustering: Hiperpersonalización"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "[INSERTAR IMAGEN 2: Propensión por clusters y productos]\n\nInsight: Particulares Premium Digitales dominan en Pension Plan (área Top)\nLeads Jóvenes: propensión mínima → DESCARTAR de campañas costosas"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ==========================================
# TAREA 3: CLUSTERING Y SEGMENTACIÓN
# ==========================================

add_section_title(prs, "TAREA 3: CLUSTERING Y SEGMENTACIÓN", "Un Mensaje Para Todos = Mensaje Para Nadie")

# Slide 13: ¿POR QUÉ SEGMENTAR?
slide = add_content_slide(
    prs,
    "¿Por Qué Segmentar?",
    [
        "El problema del marketing 'one-size-fits-all':",
        "❌ Mismo mensaje para cliente de 25 y 45 años → Nadie se siente hablado",
        "❌ Oferta genérica a premium y básico → Baja relevancia",
        "",
        "Solución: Clustering (Machine Learning No Supervisado)",
        "✅ Agrupar clientes con características similares",
        "✅ Personalizar mensajes por segmento",
        "✅ Optimizar inversión en segmentos rentables",
        "",
        "Técnica: PCA (6 componentes) + K-Means (6 clusters)",
        "Variables: Sociodemográficas, Comportamiento, Geográficas, Productos"
    ]
)

# Slide 14: LOS 6 PERFILES
slide = add_content_slide(
    prs,
    "Los 6 Perfiles de Cliente",
    [
        "1. Particulares Sin Cuenta (10.8% - 47K): Reengagement, bajo valor",
        "",
        "2. Particulares Premium Digitales (21.9% - 97K): 🏆 TESORO",
        "   Profesionales 25-35, alta propensión Pension+Investment",
        "",
        "3. Jóvenes Recientes (16.9% - 75K): Cross-selling rápido",
        "",
        "4. Jóvenes Potenciales (25.7% - 114K): Nurturing largo plazo",
        "",
        "5. Leads Jóvenes (14.4% - 64K): Bajo coste, esperar maduración",
        "",
        "6. Jóvenes Básicos Activos (10.3% - 46K): Alta para Investment"
    ]
)

# Slide 15: ESTRATEGIA PRIORIZACIÓN
slide = add_content_slide(
    prs,
    "Matriz de Priorización: Propensión × Cluster",
    [
        "🎯 PRIORIDAD MÁXIMA (80% presupuesto):",
        "• Cluster 2: Particulares Premium Digitales",
        "  → Propensión Top Pension+Investment (€6K-1.7K margen)",
        "• Cluster 6: Jóvenes Básicos Activos",
        "  → Propensión High Investment+Tarjeta",
        "",
        "⚠️ PRIORIDAD MEDIA (15% presupuesto):",
        "• Cluster 3: Jóvenes Recientes → Cross-selling automatizado",
        "",
        "❌ DESCARTAR (0-5% presupuesto):",
        "• Cluster 5: Leads Jóvenes → ROI negativo",
        "• Cluster 1: Particulares Sin Cuenta → Inactivos",
        "",
        "Resultado: 80% presupuesto en 32% clientes (Clusters 2+6)"
    ]
)

# Slide 16: PERSONALIZACIÓN
slide = add_content_slide(
    prs,
    "Mismo Producto, Diferente Historia",
    [
        "Ejemplo: Campaña PENSION PLAN",
        "",
        "Para 'Particulares Premium Digitales':",
        "• Mensaje: 'Optimiza tu futuro: Plan con ventajas fiscales'",
        "• Tono: Profesional, datos, beneficio fiscal",
        "• Canal: Email ejecutivo + app notification",
        "• Incentivo: Asesoría gratuita personalizada",
        "",
        "Para 'Jóvenes Básicos Activos':",
        "• Mensaje: 'Tu yo del futuro te lo agradecerá: 50€/mes'",
        "• Tono: Cercano, aspiracional, simple",
        "• Canal: Push app + Instagram ads",
        "• Incentivo: Bono bienvenida €50",
        "",
        "Resultado: Conversion rate 3-5x superior vs campaña genérica"
    ]
)

# ==========================================
# TAREA 4: VALOR CLIENTE + ESTRATEGIA
# ==========================================

add_section_title(prs, "TAREA 4: VALOR DEL CLIENTE Y ESTRATEGIA", "El ROI Lo Es Todo")

# Slide 17: VALOR DEL CLIENTE
slide = add_content_slide(
    prs,
    "No Se Trata de Vender Más, Sino de Vender Mejor",
    [
        "Ranking productos por margen:",
        "1. Pension Plan: €6,071 → 1 venta = 100 ventas tarjeta",
        "2. Investment: €1,699",
        "3. Payment Card: €60",
        "4. Account: €70",
        "",
        "Ejemplo Payment Card:",
        "• 296 conversiones × €60 = €17,760",
        "• Coste campaña: ~€30,000",
        "• ROI: NEGATIVO (-€12,240)",
        "",
        "Ejemplo Pension Plan:",
        "• 502 conversiones × €6,071 = €3,047,642",
        "• Coste campaña: ~€320,000",
        "• ROI: +€2,727,642 (851% ROI)",
        "",
        "Conclusión: Priorizar brutalmente en productos alto margen"
    ]
)

# Slide 18: ESCENARIOS PENSION PLAN
slide = add_content_slide(
    prs,
    "¿Cuánto Podemos Ganar? (Pension Plan)",
    [
        "Base: Conversión 3.8% validada, Margen €6,071/cliente",
        "",
        "ESCENARIO 1 - BÁSICO (solo Top):",
        "• 8,629 clientes → 328 conversiones → €1,991,288",
        "• Incremento vs orgánico: +149%",
        "",
        "ESCENARIO 2 - MEDIO (Top+High): ⭐ RECOMENDADO",
        "• 14,403 clientes → 547 conversiones → €3,320,837",
        "• Incremento vs orgánico: +315%",
        "",
        "ESCENARIO 3 - AMPLIO (Top+High+Medium):",
        "• 17,654 clientes → 671 conversiones → €4,073,641",
        "• Incremento vs orgánico: +409%",
        "",
        "Proyección 3 meses: €5.4M - €7.2M revenue",
        "(vs €2.4M orgánico = +€3M - €4.8M incrementales)"
    ]
)

# Slide 19: CANALES MARKETING
slide = add_content_slide(
    prs,
    "Optimización de Canales: No Todos Convierten Igual",
    [
        "Canales de ALTA conversión (Pension Plan):",
        "1. Canal KAT: 22% conversión 🏆 → 40% presupuesto",
        "2. Canal 003: 18% conversión → 25% presupuesto",
        "3. Canal 007: 15% conversión → 20% presupuesto",
        "4. Canal KFC: 12% conversión → 15% presupuesto",
        "",
        "Resto de canales: Conversión <5% → ELIMINAR",
        "",
        "Resultado concentración:",
        "• Conversión global: 16.75% (vs 6% actual)",
        "• 14,403 clientes × 16.75% = 2,413 conversiones",
        "• Revenue: €14.6M (vs €3.3M sin optimizar)",
        "",
        "Optimización de canales puede 4x el revenue sin aumentar budget"
    ]
)

# Slide 20: INVESTMENT Y PAYMENT CARD
slide = add_content_slide(
    prs,
    "¿Qué Hacer con Investment y Payment Card?",
    [
        "INVESTMENT FAMILY:",
        "• Conversión: 0.19% (muy baja)",
        "• Revenue esperado: €40K (escenario Top)",
        "• ROI: Apenas positivo (+€15K, 60%)",
        "• Recomendación: NO priorizar. Reevaluar producto/oferta",
        "",
        "PAYMENT CARD:",
        "• Conversión: 8.9% (alta) pero margen bajo €60",
        "• Revenue: €832K (escenario amplio)",
        "• ROI: +€682K (455%) → Positivo pero limitado",
        "• Recomendación: Estrategia secundaria (20% budget)",
        "• Uso: Producto 'anzuelo' para engagement y cross-sell posterior",
        "",
        "Conclusión: 60% esfuerzo Pension, 20% Payment Card, 20% infra"
    ]
)

# Slide 21: ESTRATEGIA 90 DÍAS
slide = add_content_slide(
    prs,
    "Plan de Acción: 3 Meses para Rentabilidad",
    [
        "MES 1 - QUICK WINS:",
        "• Lanzar Pension Plan Escenario Medio (14,403 clientes Top+High)",
        "• Clusters: 70% Premium Digitales + 30% Jóvenes Activos",
        "• Canales: KAT 40%, 003 30%, 007 30%",
        "• Objetivo: 500+ conversiones, €3M revenue",
        "",
        "MES 2 - ESCALADO:",
        "• Expandir a Escenario Amplio si Mes 1 exitoso (+3,251 Medium)",
        "• Lanzar Payment Card en Clusters 3+4 (jóvenes)",
        "• Objetivo: 750 conversiones acumuladas, €4.5M revenue",
        "",
        "MES 3 - OPTIMIZACIÓN:",
        "• Analizar datos, refinar modelos",
        "• Cross-sell Pension a compradores Payment Card (conversión 15%)",
        "• Establecer proceso repetible mensual",
        "",
        "TOTAL 3 MESES: €5.4M-€7.2M (vs €2.4M orgánico)"
    ]
)

# Slide 22: ENTREGABLES
slide = add_content_slide(
    prs,
    "Entregables Para Marketing - Listo Para Usar",
    [
        "📊 Ficheros de Datos Scored:",
        "• df_pension_plan_scored.csv (423K clientes con propensity_score)",
        "• df_payment_card_scored.csv (393K clientes)",
        "• df_inversion_family_scored.csv (431K clientes)",
        "• Importar directo a CRM/Email platform",
        "",
        "📈 Dashboard PowerBI Interactivo:",
        "• Monitorización tiempo real, Filtros por cluster/propensión/geo",
        "",
        "📝 Guía de Mensajes Personalizados:",
        "• Templates por Cluster × Producto",
        "",
        "🗓️ Calendario de Campañas:",
        "• Timing óptimo, Plan A/B testing",
        "",
        "⚙️ Soporte Técnico:",
        "• APIs scoring tiempo real, Reentrenamiento trimestral"
    ]
)

# Slide 23: CONCLUSIONES
slide = add_content_slide(
    prs,
    "De la Crisis a la Rentabilidad en 90 Días",
    [
        "LO QUE HEMOS DEMOSTRADO:",
        "",
        "✅ Datos como activo: 442K clientes = €50M+ revenue no realizado",
        "✅ Priorización: Pension (€6K) debe ser 80% del esfuerzo",
        "✅ Segmentación: 80% presupuesto en 32% clientes (Clusters 2+6)",
        "✅ Propensión: 3.8% vs 0.17% orgánico = 22x mejora validada",
        "✅ Canales: 4 estrella con conversión 12-22%, eliminar resto",
        "",
        "IMPACTO ESPERADO:",
        "• Revenue incremental: €3M - €4.8M en 90 días",
        "• ROI: 2,500% - 6,000% (Pension Plan)",
        "• Path to profitability: EBITDA positivo Q3 2019",
        "",
        "NEXT STEPS:",
        "• Semana 1: Aprobación plan | Semana 2-3: Preparación",
        "• Día 22: LANZAMIENTO Fase 1"
    ]
)

# Slide 24: CIERRE
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
title_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf = title_shape.text_frame
tf.text = "¿PREGUNTAS?"
tf.paragraphs[0].font.size = Pt(60)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
tf = subtitle.text_frame
tf.text = "Grupo 1 - Nuclio Digital School | dscesp-0525"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Guardar presentación
output_file = "/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/Presentacion_TFM_easyMoney.pptx"
prs.save(output_file)
print(f"✅ Presentación creada: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
