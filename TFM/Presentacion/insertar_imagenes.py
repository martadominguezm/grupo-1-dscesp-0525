#!/usr/bin/env python3
"""
Script para insertar las 5 imágenes PowerBI en la presentación
"""

from pptx import Presentation
from pptx.util import Inches
import os

# Rutas
pptx_file = "/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/Presentacion_TFM_easyMoney.pptx"
img_folder = "/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/imagenes_powerbi"

# Verificar que existan las imágenes
imagenes_requeridas = [
    "imagen_1_margenes_distribucion.png",
    "imagen_2_propension_clusters.png",
    "imagen_3_temporal_geografico.png",
    "imagen_4_demografico.png",
    "imagen_5_margen_familias.png"
]

print("🔍 Verificando imágenes...")
imagenes_encontradas = []
imagenes_faltantes = []

for img in imagenes_requeridas:
    img_path = os.path.join(img_folder, img)
    if os.path.exists(img_path):
        imagenes_encontradas.append(img)
        print(f"  ✅ {img}")
    else:
        imagenes_faltantes.append(img)
        print(f"  ❌ {img} - NO ENCONTRADA")

if imagenes_faltantes:
    print(f"\n⚠️  Faltan {len(imagenes_faltantes)} imágenes por guardar:")
    for img in imagenes_faltantes:
        print(f"   - {img}")
    print("\nPor favor, guarda las imágenes faltantes según las instrucciones")
    print("y vuelve a ejecutar este script.")
    exit(1)

print(f"\n✅ Todas las imágenes encontradas ({len(imagenes_encontradas)}/5)")
print("\n📊 Insertando imágenes en PowerPoint...\n")

# Abrir presentación
prs = Presentation(pptx_file)

def eliminar_placeholders_grises(slide):
    """Elimina los textboxes con placeholders de imágenes"""
    shapes_to_delete = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "[INSERTAR IMAGEN" in text or "IMAGEN POWERBI" in text.upper():
                shapes_to_delete.append(shape)

    # Eliminar en orden inverso para no afectar índices
    for shape in reversed(shapes_to_delete):
        sp = shape.element
        sp.getparent().remove(sp)

def insertar_imagen(slide, img_path, left, top, width, height):
    """Inserta imagen en slide"""
    # Eliminar placeholders grises primero
    eliminar_placeholders_grises(slide)

    # Insertar imagen
    slide.shapes.add_picture(
        img_path,
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height)
    )

# Mapeo: Slide index → (imagen, left, top, width, height)
# Los slides están indexados desde 0
inserciones = {
    3: ("imagen_3_temporal_geografico.png", 4, 1.5, 5.5, 5),      # Slide 4 (index 3)
    4: ("imagen_4_demografico.png", 4, 1.5, 5.5, 5),              # Slide 5 (index 4)
    5: ("imagen_1_margenes_distribucion.png", 4.5, 1.5, 5, 5),    # Slide 6 (index 5)
    6: ("imagen_5_margen_familias.png", 4.5, 1.5, 5, 5),          # Slide 7 (index 6)
    12: ("imagen_2_propension_clusters.png", 0.5, 1.5, 9, 5),     # Slide 13 (index 12)
}

# Insertar imágenes
for slide_idx, (img_name, left, top, width, height) in inserciones.items():
    try:
        slide = prs.slides[slide_idx]
        img_path = os.path.join(img_folder, img_name)

        insertar_imagen(slide, img_path, left, top, width, height)

        # Obtener título del slide
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        print(f"✅ Slide {slide_idx + 1}: {img_name}")
        if title:
            print(f"   '{title}'")

    except Exception as e:
        print(f"❌ Error en slide {slide_idx + 1}: {e}")

# Guardar presentación actualizada
output_file = "/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/Presentacion_TFM_easyMoney_CON_IMAGENES.pptx"
prs.save(output_file)

print(f"\n🎉 ¡COMPLETADO!")
print(f"📊 Presentación guardada: Presentacion_TFM_easyMoney_CON_IMAGENES.pptx")
print(f"📁 Ubicación: {output_file}")
print(f"\n✨ La presentación ahora tiene las 5 imágenes PowerBI insertadas")
print(f"   Abre el archivo y verifica que todo se vea bien.")
