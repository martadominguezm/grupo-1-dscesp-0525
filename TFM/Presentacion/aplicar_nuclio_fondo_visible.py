#!/usr/bin/env python3
"""
Aplicar colores Nuclio con FONDO MUY VISIBLE (no gris claro)
Usar crema/beige cálido que contraste más con blanco
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# COLORES NUCLIO
NUCLIO = {
    'naranja': RGBColor(255, 107, 53),        # #FF6B35 - Naranja Nuclio
    'naranja_fondo': RGBColor(255, 200, 150), # #FFC896 - Naranja claro para fondo
    'negro': RGBColor(26, 26, 26),            # #1A1A1A - Negro
    'blanco': RGBColor(255, 255, 255),        # #FFFFFF - Blanco
}

def add_background(slide, color):
    """Añade fondo de color usando método exitoso"""
    fondo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = color
    fondo.line.fill.background()

    # Mover al fondo
    slide.shapes._spTree.remove(fondo._element)
    slide.shapes._spTree.insert(0, fondo._element)
    return fondo

def add_bar(slide, color, position='top'):
    """Añade barra decorativa"""
    if position == 'top':
        left, top, width, height = 0, 0, 10, 0.2
    else:
        left, top, width, height = 0, 7.3, 10, 0.2

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

print("🎨 Aplicando colores Nuclio con FONDO NARANJA...\n")

# Abrir la presentación CON IMÁGENES
prs = Presentation("/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/Presentacion_TFM_easyMoney_CON_IMAGENES.pptx")

# USAR FONDO NARANJA CLARO (como en el test que funcionó)
FONDO_ELEGIDO = NUCLIO['naranja_fondo']  # Naranja claro VISIBLE

print(f"🎨 Color de fondo: Naranja claro (255, 200, 150) - VISIBLE\n")

for idx, slide in enumerate(prs.slides):
    slide_num = idx + 1

    # FONDO BEIGE VISIBLE
    add_background(slide, FONDO_ELEGIDO)

    # BARRAS según slide (alternando naranja y negro)
    if idx == 0:  # Portada
        add_bar(slide, NUCLIO['naranja'], 'top')
        add_bar(slide, NUCLIO['negro'], 'bottom')
    elif 1 <= idx <= 6:  # Tarea 1
        add_bar(slide, NUCLIO['naranja'], 'top')
    elif idx == 7:  # Separador Tarea 2
        add_bar(slide, NUCLIO['negro'], 'top')
        add_bar(slide, NUCLIO['negro'], 'bottom')
    elif 8 <= idx <= 12:  # Tarea 2
        add_bar(slide, NUCLIO['negro'], 'top')
    elif idx == 13:  # Separador Tarea 3
        add_bar(slide, NUCLIO['naranja'], 'top')
        add_bar(slide, NUCLIO['naranja'], 'bottom')
    elif 14 <= idx <= 17:  # Tarea 3
        add_bar(slide, NUCLIO['naranja'], 'top')
    elif idx == 18:  # Separador Tarea 4
        add_bar(slide, NUCLIO['negro'], 'top')
        add_bar(slide, NUCLIO['negro'], 'bottom')
    elif 19 <= idx <= 25:  # Tarea 4
        add_bar(slide, NUCLIO['negro'], 'top')
    elif idx == 26:  # Cierre
        add_bar(slide, NUCLIO['naranja'], 'top')
        add_bar(slide, NUCLIO['negro'], 'bottom')

    # Estilizar títulos
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        for p in tf.paragraphs:
            if p.text:
                p.font.bold = True

                if idx == 0:  # Portada
                    p.font.size = Pt(44)
                    p.font.color.rgb = NUCLIO['negro']
                elif idx in [7, 13, 18]:  # Separadores
                    p.font.size = Pt(48)
                    p.font.color.rgb = NUCLIO['blanco']  # Blanco sobre color
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.size = Pt(36)
                    if 1 <= idx <= 6:
                        p.font.color.rgb = NUCLIO['naranja']
                    elif 8 <= idx <= 12:
                        p.font.color.rgb = NUCLIO['negro']
                    elif 14 <= idx <= 17:
                        p.font.color.rgb = NUCLIO['naranja']
                    else:
                        p.font.color.rgb = NUCLIO['negro']

    # Estilizar contenido
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            for p in tf.paragraphs:
                if p.text:
                    p.font.color.rgb = NUCLIO['negro']
                    if p.level == 0:
                        p.font.size = Pt(18)
                    else:
                        p.font.size = Pt(16)

    print(f"✅ Slide {slide_num}: Fondo naranja + barras Nuclio aplicados")

# Guardar
output = "/Users/alex/Desktop/TFM_final/grupo-1-dscesp-0525/TFM/Presentacion/Presentacion_TFM_NUCLIO_FONDO_NARANJA.pptx"
prs.save(output)

print(f"\n🎉 ¡Presentación con fondo NARANJA creada!")
print(f"📊 Archivo: Presentacion_TFM_NUCLIO_FONDO_NARANJA.pptx")
print(f"\n🎨 Paleta Nuclio aplicada:")
print(f"   • Fondo: Naranja claro (255, 200, 150) - VISIBLE")
print(f"   • Barras decorativas: Naranja oscuro + Negro alternando")
print(f"   • Textos: Negro sobre fondo naranja")
print(f"   • Títulos separadores: Blanco sobre barras de color")
print(f"\n🍊 Ahora SÍ deberías ver el fondo naranja claramente")
